"""文档加载器 - 解析多种文档格式提取纯文本"""

from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


class DocumentLoader:
    """
    文档加载器

    支持格式：txt, md, pdf, docx, pptx, xlsx, html
    返回纯文本字符串，供后续分块和向量化。
    """

    SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx", ".xlsx", ".html", ".htm"}

    def load(self, file_path: str) -> str:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        ext = path.suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"不支持的文件格式: {ext}，支持: {self.SUPPORTED_EXTENSIONS}")

        if ext in (".txt", ".md"):
            return self._load_text(path)
        elif ext == ".pdf":
            return self._load_pdf(path)
        elif ext == ".docx":
            return self._load_docx(path)
        elif ext == ".pptx":
            return self._load_pptx(path)
        elif ext == ".xlsx":
            return self._load_xlsx(path)
        elif ext in (".html", ".htm"):
            return self._load_html(path)

    def load_bytes(self, content: bytes, filename: str) -> str:
        """从字节流加载，用于文件上传场景"""
        ext = Path(filename).suffix.lower()
        if ext in (".txt", ".md"):
            return content.decode("utf-8", errors="ignore")
        elif ext == ".pdf":
            return self._load_pdf_bytes(content)
        elif ext == ".docx":
            return self._load_docx_bytes(content)
        elif ext == ".pptx":
            return self._load_pptx_bytes(content)
        elif ext == ".xlsx":
            return self._load_xlsx_bytes(content)
        elif ext in (".html", ".htm"):
            return self._load_html_bytes(content)
        else:
            raise ValueError(f"不支持的文件格式: {ext}")

    def _load_text(self, path: Path) -> str:
        return path.read_text(encoding="utf-8", errors="ignore")

    def _load_pdf(self, path: Path) -> str:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(str(path))
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return "\n\n".join(pages)
        except ImportError:
            logger.warning("PyPDF2 未安装，无法解析 PDF")
            raise

    def _load_pdf_bytes(self, content: bytes) -> str:
        try:
            from PyPDF2 import PdfReader
            import io
            reader = PdfReader(io.BytesIO(content))
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return "\n\n".join(pages)
        except ImportError:
            logger.warning("PyPDF2 未安装，无法解析 PDF")
            raise

    def _load_docx(self, path: Path) -> str:
        try:
            from docx import Document
            doc = Document(str(path))
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            logger.warning("python-docx 未安装，无法解析 DOCX")
            raise

    def _load_docx_bytes(self, content: bytes) -> str:
        try:
            from docx import Document
            import io
            doc = Document(io.BytesIO(content))
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            logger.warning("python-docx 未安装，无法解析 DOCX")
            raise

    def _load_pptx(self, path: Path) -> str:
        """解析 PowerPoint 文件"""
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            return self._extract_pptx_text(prs)
        except ImportError:
            logger.warning("python-pptx 未安装，无法解析 PPTX")
            raise

    def _load_pptx_bytes(self, content: bytes) -> str:
        """从字节流解析 PowerPoint"""
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(content))
            return self._extract_pptx_text(prs)
        except ImportError:
            logger.warning("python-pptx 未安装，无法解析 PPTX")
            raise

    def _extract_pptx_text(self, prs) -> str:
        """提取 PPTX 中的文本"""
        slides_text = []
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                # 提取表格内容
                if shape.has_table:
                    for row in shape.table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_texts:
                            slide_texts.append(" | ".join(row_texts))
            if slide_texts:
                slides_text.append("\n".join(slide_texts))
        return "\n\n".join(slides_text)

    def _load_xlsx(self, path: Path) -> str:
        """解析 Excel 文件"""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True, data_only=True)
            return self._extract_xlsx_text(wb)
        except ImportError:
            logger.warning("openpyxl 未安装，无法解析 XLSX")
            raise

    def _load_xlsx_bytes(self, content: bytes) -> str:
        """从字节流解析 Excel"""
        try:
            from openpyxl import load_workbook
            import io
            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            return self._extract_xlsx_text(wb)
        except ImportError:
            logger.warning("openpyxl 未安装，无法解析 XLSX")
            raise

    def _extract_xlsx_text(self, wb) -> str:
        """提取 XLSX 中的文本"""
        sheets_text = []
        for sheet in wb.worksheets:
            sheet_rows = []
            # 提取表头
            headers = []
            for cell in next(sheet.iter_rows(min_row=1, max_row=1)):
                val = str(cell.value).strip() if cell.value is not None else ""
                if val:
                    headers.append(val)
            if headers:
                sheet_rows.append(f"[{sheet.title}] " + " | ".join(headers))

            # 提取数据行（跳过表头）
            for row in sheet.iter_rows(min_row=2, values_only=True):
                row_texts = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if row_texts:
                    sheet_rows.append(" | ".join(row_texts))

            if sheet_rows:
                sheets_text.append("\n".join(sheet_rows))
        return "\n\n".join(sheets_text)

    def _load_html(self, path: Path) -> str:
        """解析 HTML 文件"""
        try:
            from bs4 import BeautifulSoup
            html_content = path.read_text(encoding="utf-8", errors="ignore")
            return self._extract_html_text(html_content)
        except ImportError:
            logger.warning("beautifulsoup4 未安装，无法解析 HTML")
            raise

    def _load_html_bytes(self, content: bytes) -> str:
        """从字节流解析 HTML"""
        try:
            from bs4 import BeautifulSoup
            html_content = content.decode("utf-8", errors="ignore")
            return self._extract_html_text(html_content)
        except ImportError:
            logger.warning("beautifulsoup4 未安装，无法解析 HTML")
            raise

    def _extract_html_text(self, html: str) -> str:
        """提取 HTML 中的文本"""
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html, "lxml")

        # 移除 script 和 style 标签
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        # 提取文本
        text = soup.get_text(separator="\n", strip=True)

        # 清理多余空行
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)
